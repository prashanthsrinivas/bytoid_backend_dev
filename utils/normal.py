import os
from pathlib import Path
import yaml
import re
import json
from datetime import datetime, timedelta
from docx import Document
from pptx import Presentation
import pytz
from bs4 import BeautifulSoup

# from astral import LocationInfo
# from astral.sun import sun


def ensure_dir(path):
    os.makedirs(path, exist_ok=True)


def load_yaml_file(path):
    if not os.path.exists(path):
        return None  # Important: use None to differentiate from empty list
    try:
        with open(path, "r", encoding="utf-8") as f:
            return yaml.safe_load(f) or []
    except yaml.YAMLError as e:
        print(f"❌ Error reading YAML file at {path}: {e}")
        return []


def can_reply_to_email(email: str) -> bool:
    """
    Returns True if it is safe to reply to this email.

    Rules:
    1. Skip system/no-reply/admin/bounce/automation emails.
    2. Skip emails from blocked domains.
    3. Only allow .com, .in, .ai TLDs.
    4. EXCEPTION: any @bytoid.<anything> is always allowed.
    """
    if not email:
        return False

    email = email.lower().strip()

    # Extract local part + domain
    match = re.match(r"^[\w\.-]+@([\w\.-]+)$", email)
    if not match:
        return False

    full_domain = match.group(1)  # example: bytoid.ca
    parts = full_domain.split(".")
    if len(parts) < 2:
        return False

    domain_name = parts[0]  # bytoid
    tld = parts[-1]  # ca

    # ✅ 4️⃣ BYPASS RULE — ALWAYS ALLOW bytoid.*
    if domain_name == "bytoid":
        return True

    # 1️⃣ Blocked keywords in local part
    blocked_keywords = [
        "no-reply",
        "noreply",
        "donotreply",
        "do-not-reply",
        "system",
        "postmaster",
        "bounce",
        "mailer-daemon",
        "undeliverable",
        "return",
        "notifications",
        "alerts",
        "updates",
        "robot",
        "automation",
    ]
    if any(k in email for k in blocked_keywords):
        return False

    # 2️⃣ Blocked domains
    blocked_domains = {
        "naukri",
        "google",
        "amazon",
        "twitter",
        "x",
        "indeed",
        "train",
        "zomato",
        "swiggy",
        "ola",
        "rapido",
        "linkedin",
        "microsoft",
        "facebook",
        "instagram",
        "youtube",
        "flipkart",
        "ubereats",
        "booking",
        "airbnb",
        "phonepe",
        "upstox",
        "irctc",
        "aubank",
        "kotak",
        "railone",
    }

    # 2️⃣ Block domain if in blocked list
    if domain_name in blocked_domains:
        return False

    # 3️⃣ Only allow TLDs: .com .in .ai
    if tld not in {"com", "in", "ai"}:
        return False

    return True


def read_function_jsons():
    """
    Reads all function JSON files and returns a string formatted for the prompt,
    including function name, description, and arguments.
    """
    paths = [
        "playbook/fn_configs/automate_functions.json",
        "playbook/fn_configs/gmail_functions.json",
        "playbook/fn_configs/google_meet_functions.json",
        "playbook/fn_configs/umail_auto_functions.json",
        "playbook/fn_configs/outlook_functions.json",
        # "playbook/fn_configs/twillo_fucntion.json",
    ]
    all_functions_details = []

    for path in paths:
        if os.path.exists(path):
            try:
                with open(path, "r", encoding="utf-8") as f:
                    data = json.load(f)
            except json.JSONDecodeError:
                continue

            for fn in data.get("functions", []):
                fn_name = fn.get("name", "Unknown Function")
                fn_description = fn.get("description", "No description provided.")
                fn_args = fn.get("args", {})
                fn_notes = fn.get("notes", {})
                fn_ret = fn.get("returns", {})

                args_list = [
                    f"    - {arg_name}: {arg_desc}"
                    for arg_name, arg_desc in fn_args.items()
                ]
                return_list = [
                    f"    - {ret_name}: {ret_desc}"
                    for ret_name, ret_desc in fn_ret.items()
                ]
                args_str = "\n".join(args_list) or "    - None"
                ret_args_str = "\n".join(return_list) or "    - None"

                function_detail = (
                    f"### Function: **{fn_name}**\n"
                    f"Description: {fn_description}\n"
                    f"Arguments:\n"
                    f"{args_str} \n"
                    f"{fn_notes} \n"
                    f"returns:{ret_args_str}"
                )
                all_functions_details.append(function_detail)

    return "\n\n".join(all_functions_details)


# print(read_function_jsons())
def read_function_jsons2(Full=False):
    """
    Reads all function JSON files and returns a dictionary of functions.

    If Full=False (default):
        returns { function_name: description }

    If Full=True:
        returns { function_name: full_function_dict }
        where full_function_dict includes name, description, parameters, examples, etc.
    """

    paths = [
        "playbook/fn_configs/automate_functions.json",
        "playbook/fn_configs/gmail_functions.json",
        "playbook/fn_configs/google_meet_functions.json",
        "playbook/fn_configs/umail_auto_functions.json",
        "playbook/fn_configs/outlook_functions.json",
        # "playbook/fn_configs/twillo_fucntion.json",
    ]

    all_functions = {}

    for path in paths:
        if os.path.exists(path):
            try:
                with open(path, "r", encoding="utf-8") as f:
                    data = json.load(f)
            except json.JSONDecodeError:
                print(f"⚠️ Skipped invalid JSON file: {path}")
                continue

            for fn in data.get("functions", []):
                fn_name = fn.get("name")
                if not fn_name:
                    continue

                if Full:
                    # Return the entire function metadata object
                    all_functions[fn_name] = fn
                else:
                    # Return only name → description mapping
                    all_functions[fn_name] = fn.get(
                        "description", "No description provided."
                    )

    return all_functions


def parse_pptx_content_to_json(content: str) -> dict:
    slides_data = {"slides": []}
    current_slide = None

    for line in content.splitlines():
        line = line.strip()
        if not line:
            continue

        # Detect slide heading: Roman numeral or numbered heading
        if re.match(r"^[IVXLCDM]+\. ", line):  # I., II., III., etc.
            # Save previous slide
            if current_slide:
                slides_data["slides"].append(current_slide)
            current_slide = {"title": "", "bulletPoints": []}
            continue

        # Detect title line
        title_match = re.match(r"^\* Title:\s*(.*)", line)
        if title_match:
            if current_slide is None:
                current_slide = {"title": "", "bulletPoints": []}
            current_slide["title"] = title_match.group(1)
            continue

        # Detect bullet lines
        bullet_match = re.match(r"^\* (?!Title:)(.*)", line)
        if bullet_match and current_slide:
            current_slide["bulletPoints"].append(bullet_match.group(1))
            continue

    if current_slide:
        slides_data["slides"].append(current_slide)

    return slides_data


def prepare_docx_data_from_ai(ai_content: str) -> dict:
    """
    Convert AI content (raw string or JSON) into doc_data format for .docx:
    {
        "title": "...",
        "sections": [{"heading": "...", "paragraphs": ["..."]}]
    }
    """
    # Try parsing ai_content as JSON first
    try:
        parsed = json.loads(ai_content)
        if "content" in parsed:
            content_text = parsed["content"].strip()
        else:
            content_text = ai_content.strip()
    except (json.JSONDecodeError, TypeError):
        content_text = ai_content.strip()

    # Split into lines
    lines = [line.strip() for line in content_text.splitlines() if line.strip()]

    # First line as title
    title = lines[0] if lines else "Document"

    # Split rest into sections
    sections = []
    current_heading = ""
    current_paragraphs = []

    for line in lines[1:]:
        if (
            line.endswith(":")
            or line.startswith("###")
            or line.endswith("------------")
        ):
            # Treat as a heading
            if current_heading or current_paragraphs:
                sections.append(
                    {
                        "heading": current_heading or "Section",
                        "paragraphs": current_paragraphs,
                    }
                )
            current_heading = line.replace("###", "").replace("-", "").strip()
            current_paragraphs = []
        else:
            current_paragraphs.append(line)
    # Append last section
    if current_heading or current_paragraphs:
        sections.append(
            {
                "heading": current_heading or "Section",
                "paragraphs": current_paragraphs or ["Content unavailable"],
            }
        )

    return {
        "title": title,
        "sections": sections
        or [{"heading": "Section 1", "paragraphs": ["Content unavailable"]}],
    }


def save_docx_from_json(doc_data, file_path):
    """Save structured Word content as a .docx file, replacing if exists"""
    file_path = Path(file_path)
    if file_path.exists():
        file_path.unlink()  # Delete existing file
    # print("doc data", doc_data)

    doc = Document()
    doc.add_heading(doc_data.get("title", ""), level=0)
    for section in doc_data.get("sections", []):
        doc.add_heading(section.get("heading", ""), level=1)
        for para in section.get("paragraphs", []):
            doc.add_paragraph(para)
    doc.save(file_path)
    # os.remove(file_path)


def save_pptx_from_json(slide_data, file_path):
    """Save structured PowerPoint content as a .pptx file, replacing if exists"""
    file_path = Path(file_path)
    if file_path.exists():
        file_path.unlink()  # Delete existing file
    # print("slide data", slide_data)

    prs = Presentation()
    for slide_info in slide_data.get("slides", []):
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_info.get("title", "")
        content = "\n".join(slide_info.get("bulletPoints", []))
        slide.placeholders[1].text = content
    prs.save(file_path)
    # os.remove(file_path)


def convert_human_date(value, base_date=None, tz_str="Asia/Kolkata"):
    """
    Convert human-readable dates or ISO YYYY-MM-DD to a tz-aware datetime
    Default time: 09:00
    Rules:
    - If year is missing, infer next valid future date
    - Never return past dates
    """
    tz = pytz.timezone(tz_str)
    if base_date is None:
        base_date = datetime.now(tz)

    value = str(value).strip().lower()
    dt = None

    # ----------------------------
    # 1. Simple keywords
    # ----------------------------
    if value in ["today", ""]:
        dt = base_date
    elif value == "tomorrow":
        dt = base_date + timedelta(days=1)

    # ----------------------------
    # 2. ISO date
    # ----------------------------
    elif re.match(r"\d{4}-\d{2}-\d{2}$", value):
        try:
            dt = datetime.strptime(value, "%Y-%m-%d")
        except Exception:
            return None

    # ----------------------------
    # 3. Numeric month-day (12-10)
    # ----------------------------
    elif re.match(r"\d{1,2}-\d{1,2}$", value):
        month, day = map(int, value.split("-"))
        dt = base_date.replace(month=month, day=day)

    # ----------------------------
    # 4. Day-month (12-jan / 12-january / 12th january)
    # ----------------------------
    elif match := re.match(r"(\d{1,2})(st|nd|rd|th)?[\s\-]+([a-z]+)", value):
        day = int(match.group(1))
        month_str = match.group(3).capitalize()

        try:
            month = datetime.strptime(month_str[:3], "%b").month
        except ValueError:
            return None

        dt = base_date.replace(month=month, day=day)

    # ----------------------------
    # 5. Relative days
    # ----------------------------
    elif match := re.match(r"(\d+)\s+days\s+from\s+now", value):
        dt = base_date + timedelta(days=int(match.group(1)))

    # ----------------------------
    # 6. Next weekday
    # ----------------------------
    elif match := re.match(r"next\s+(\w+)", value):
        weekdays = [
            "monday",
            "tuesday",
            "wednesday",
            "thursday",
            "friday",
            "saturday",
            "sunday",
        ]
        day_name = match.group(1).lower()
        if day_name not in weekdays:
            return None

        target_wd = weekdays.index(day_name)
        delta = (target_wd - base_date.weekday() + 7) or 7
        dt = base_date + timedelta(days=delta)

    else:
        return None

    # ----------------------------
    # 7. Infer year (never backward)
    # ----------------------------
    if dt.year == base_date.year:
        if dt.date() < base_date.date():
            dt = dt.replace(year=base_date.year + 1)

    # ----------------------------
    # 8. Default time 09:00
    # ----------------------------
    dt = dt.replace(hour=9, minute=0, second=0, microsecond=0)

    # ----------------------------
    # 9. Ensure tz-aware
    # ----------------------------
    if dt.tzinfo is None:
        dt = tz.localize(dt)
    else:
        dt = dt.astimezone(tz)

    return dt


# ------------------------
# Human-readable Time
# ------------------------
def convert_human_time(value, base_date=None, tz_str="Asia/Kolkata"):
    """
    Convert human-readable time expressions to a tz-aware datetime.

    Supported:
    - 10, 10am, 10 am, 10:30, 10:30 pm
    - 10-11 / 10 to 11 (returns start time)
    - quarter past 10, half past 10, quarter to 11
    - now, in 30 minutes, after 1 hour
    - morning, afternoon, evening, night
    - noon, midnight
    """
    tz = pytz.timezone(tz_str)
    if base_date is None:
        base_date = datetime.now(tz)

    value = str(value).strip().lower()
    hour, minute = None, None

    # ----------------------------
    # 1. Immediate keywords
    # ----------------------------
    if value in ["now", "right now"]:
        dt = base_date

    elif m := re.match(r"in (\d+) minutes?", value):
        dt = base_date + timedelta(minutes=int(m.group(1)))

    elif m := re.match(r"(in|after) (\d+) hours?", value):
        dt = base_date + timedelta(hours=int(m.group(2)))

    # ----------------------------
    # 2. Time ranges (10-11 / 10 to 11)
    # ----------------------------
    elif m := re.match(r"(\d{1,2})(:\d{2})?\s*(am|pm)?\s*(to|\-)\s*(\d{1,2})", value):
        hour = int(m.group(1))
        minute = int(m.group(2)[1:]) if m.group(2) else 0
        ampm = m.group(3)

        if ampm == "pm" and hour != 12:
            hour += 12
        elif ampm == "am" and hour == 12:
            hour = 0

        dt = base_date.replace(hour=hour, minute=minute)

    # ----------------------------
    # 3. HH:MM AM/PM
    # ----------------------------
    elif m := re.match(r"(\d{1,2}):(\d{2})\s*(am|pm)?", value):
        hour, minute = int(m.group(1)), int(m.group(2))
        if m.group(3) == "pm" and hour != 12:
            hour += 12
        elif m.group(3) == "am" and hour == 12:
            hour = 0
        dt = base_date.replace(hour=hour, minute=minute)

    # ----------------------------
    # 4. 10 am / 10pm
    # ----------------------------
    elif m := re.match(r"(\d{1,2})\s*(am|pm)", value):
        hour = int(m.group(1))
        minute = 0
        if m.group(2) == "pm" and hour != 12:
            hour += 12
        elif m.group(2) == "am" and hour == 12:
            hour = 0
        dt = base_date.replace(hour=hour, minute=minute)

    # ----------------------------
    # 5. Plain number (10 → 10:00)
    # ----------------------------
    elif value.isdigit():
        hour = int(value)
        minute = 0
        dt = base_date.replace(hour=hour, minute=minute)

    # ----------------------------
    # 6. Natural language clock phrases
    # ----------------------------
    elif m := re.match(r"quarter to (\d+)", value):
        hour = int(m.group(1)) - 1
        minute = 45
        dt = base_date.replace(hour=hour, minute=minute)

    elif m := re.match(r"quarter past (\d+)", value):
        hour = int(m.group(1))
        minute = 15
        dt = base_date.replace(hour=hour, minute=minute)

    elif m := re.match(r"half past (\d+)", value):
        hour = int(m.group(1))
        minute = 30
        dt = base_date.replace(hour=hour, minute=minute)

    elif m := re.match(r"(\d+) past (\d+)", value):
        minute, hour = int(m.group(1)), int(m.group(2))
        dt = base_date.replace(hour=hour, minute=minute)

    elif m := re.match(r"(\d+) to (\d+)", value):
        minute = 60 - int(m.group(1))
        hour = int(m.group(2)) - 1
        dt = base_date.replace(hour=hour, minute=minute)

    # ----------------------------
    # 7. Named time buckets
    # ----------------------------
    elif value in ["morning"]:
        dt = base_date.replace(hour=9, minute=0)
    elif value in ["afternoon"]:
        dt = base_date.replace(hour=14, minute=0)
    elif value in ["evening"]:
        dt = base_date.replace(hour=18, minute=0)
    elif value in ["night"]:
        dt = base_date.replace(hour=21, minute=0)
    elif value in ["noon", "lunch"]:
        dt = base_date.replace(hour=12, minute=0)
    elif value in ["midnight"]:
        dt = base_date.replace(hour=0, minute=0)

    else:
        return None

    # ----------------------------
    # 8. Normalize seconds + tz
    # ----------------------------
    dt = dt.replace(second=0, microsecond=0)

    if dt.tzinfo is None:
        dt = tz.localize(dt)
    else:
        dt = dt.astimezone(tz)

    return dt


EMAIL_TITLES = [
    "AI Automation",
    "SaaS Growth Strategies",
    "Developer Tools Evolution",
    "Modern Cybersecurity",
    "Cloud Computing Trends",
    "Workflow Orchestration",
    "Customer Engagement Intelligence",
    "Productivity Hacks 2025",
    "Edge Computing",
    "Microservice Architecture",
    "Zero Trust Security",
    "Smart Assistants",
    "Intelligent Workflows",
    "Data Privacy Essentials",
    "Machine Learning at Work",
    "API-First Infrastructure",
    "Serverless Technology",
    "Quantum Computing",
    "Observability & Monitoring",
    "Cloud-Native Engineering",
    "DevOps Culture",
    "Platform Engineering",
    "LLM-Powered Applications",
    "Automation Pipelines",
    "Adaptive Workflows",
    "ChatOps Workflows",
    "Event-Driven Systems",
    "Composable SaaS",
    "Realtime Analytics",
    "AI for Developers",
    "Security Automation",
    "Email Automation",
    "Smart Scheduling",
    "AI-Powered Customer Support",
    "Autonomous Agents",
    "Identity Management",
    "API Integration",
    "Container Orchestration",
    "Kubernetes Automation",
    "Tech Workflow Insights",
    "Business Automation",
    "Digital Transformation",
    "AI-Augmented Workforce",
    "Intelligent Routing",
    "Smart CRM Automation",
    "Cloud Optimization",
    "Future of Work",
    "AI Productivity Boost",
    "Scalable Infrastructure",
    "AI-Driven Work Management",
]


# --------------------------
# SUBJECT EXTRACTOR
# --------------------------


def extract_subject_from_html(html: str, fallback: str) -> str:
    try:
        soup = BeautifulSoup(html, "html.parser")

        # 1) Prefer <title>
        title_tag = soup.find("title")
        if title_tag and title_tag.text.strip():
            return title_tag.text.strip()

        # 2) Or an <h1>
        h1_tag = soup.find("h1")
        if h1_tag and h1_tag.text.strip():
            return h1_tag.text.strip()

    except Exception:
        pass

    # fallback
    return fallback


import re


def remove_not_found_entities(text: str, not_found: list[str]) -> str:
    if not text or not not_found:
        return text

    cleaned = text

    for name in not_found:
        escaped = re.escape(name)

        # remove patterns like:
        # ", Josh"
        # "and Josh"
        # "& Josh"
        # "Josh,"
        patterns = [
            rf"\s*,\s*{escaped}\b",
            rf"\b{escaped}\s*,\s*",
            rf"\s+(and|&)\s+{escaped}\b",
            rf"\b{escaped}\s+(and|&)\s+",
            rf"\b{escaped}\b",
        ]

        for pattern in patterns:
            cleaned = re.sub(pattern, " ", cleaned, flags=re.IGNORECASE)

    # normalize whitespace
    cleaned = re.sub(r"\s{2,}", " ", cleaned).strip()

    # remove dangling "with", "and"
    cleaned = re.sub(r"\b(with|and)\s*$", "", cleaned, flags=re.IGNORECASE).strip()

    return cleaned
